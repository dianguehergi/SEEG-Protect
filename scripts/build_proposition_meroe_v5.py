"""Create the MEROE V5 scoring workbook and V4 executive pitch."""

from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
XLSX = DOCS / "PROPOSITION_MEROE_DG_V5.xlsx"
PPTX = DOCS / "PROPOSITION_MEROE_PITCH_V4.pptx"
BACKGROUND = ROOT / "powerbi_meroe_v312" / "assets" / "meroe_dashboard_background_v1.png"

NAVY, PANEL, GREEN, BLUE, GOLD, RED, WHITE, MUTED = (
    "071521", "102B3D", "20D6B5", "5B8CFF", "F5B942", "FF6B6B", "F4F7FA", "8DA2B5"
)

SCORE_ROWS = [
    ("B1", "Données abonné", 20, "Puissance souscrite > 15 KVA en domestique", 10),
    ("B1", "Données abonné", 20, "Changement puissance < 3 mois", 5),
    ("B1", "Données abonné", 20, "Adresse différente du quartier SAP", 5),
    ("B2", "Historique 6 mois", 30, "Baisse > 80 % vs moyenne 6 mois", 15),
    ("B2", "Historique 6 mois", 30, "Baisse 60 % à 80 %", 12),
    ("B2", "Historique 6 mois", 30, "Baisse 40 % à 59 %", 10),
    ("B2", "Historique 6 mois", 30, "Baisse 25 % à 39 %", 7),
    ("B2", "Historique 6 mois", 30, "Consommation nulle pendant 2 mois", 15),
    ("B2", "Historique 6 mois", 30, "Consommation nulle pendant 1 mois", 8),
    ("B2", "Historique 6 mois", 30, "Pic > 200 % vs moyenne", 10),
    ("B2", "Historique 6 mois", 30, "Écart-type > 40 %", 5),
    ("B3", "Recharge / solde", 20, "Recharge < 500 F et consommation > 10 kWh/j", 10),
    ("B3", "Recharge / solde", 20, "Fréquence de recharge doublée en 1 mois", 5),
    ("B3", "Recharge / solde", 20, "Solde négatif > 7 jours", 5),
    ("B4", "Données techniques", 20, "Ouverture cache-bornes", 10),
    ("B4", "Données techniques", 20, "Tension < 200 V ou > 250 V", 10),
    ("B4", "Données techniques", 20, "Inversion courant > 3 fois", 5),
    ("B4", "Données techniques", 20, "0 kWh en heures pleines pendant 30 jours", 5),
    ("B5", "Géolocalisation", 10, "Au moins 5 voisins en anomalie", 10),
    ("B5", "Géolocalisation", 10, "3 à 4 voisins en anomalie", 7),
    ("B5", "Géolocalisation", 10, "1 à 2 voisins en anomalie", 5),
    ("B5", "Géolocalisation", 10, "GPS différent de l’adresse de facturation", 5),
]


def style_sheet(ws, widths):
    ws.freeze_panes = "A2"
    ws.sheet_view.showGridLines = False
    for cell in ws[1]:
        cell.fill = PatternFill("solid", fgColor=NAVY)
        cell.font = Font(color=WHITE, bold=True)
        cell.alignment = Alignment(vertical="center")
    ws.row_dimensions[1].height = 28
    for idx, width in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(idx)].width = width
    thin = Side(style="thin", color="D7E1E8")
    for row in ws.iter_rows(min_row=2):
        for cell in row:
            cell.border = Border(bottom=thin)
            cell.alignment = Alignment(vertical="top", wrap_text=True)


def build_workbook():
    wb = Workbook()
    ws = wb.active
    ws.title = "Barème IA V5"
    ws.append(["Bloc", "Famille", "Plafond bloc", "Règle explicable", "Points"])
    for row in SCORE_ROWS:
        ws.append(row)
    style_sheet(ws, [10, 25, 15, 62, 12])
    for row in range(2, ws.max_row + 1):
        ws.cell(row, 1).fill = PatternFill("solid", fgColor="E8F8F4")
    decision = wb.create_sheet("Décision")
    decision.append(["Score", "Couleur", "Décision IA", "Action autorisée"])
    for row in [
        ("0–30", "VERT", "RAS", "Surveillance normale"),
        ("31–59", "ORANGE", "WATCHLIST", "Réévaluation à J+7"),
        ("60–79", "ROUGE", "ALERTE", "Contrôle terrain prioritaire"),
        ("80–100", "NOIR", "URGENT", "Contrôle urgent ; décision SEEG avant toute coupure"),
    ]:
        decision.append(row)
    style_sheet(decision, [16, 16, 22, 58])
    access = wb.create_sheet("Matrice accès")
    access.append(["Fonction", "Admin fondateur", "CTO", "Collaborateur"])
    rows = [
        ("Trésor et montants", "Complet", "Interdit", "Interdit"),
        ("PROTEC / SMS / appels", "Complet", "Complet", "Lecture limitée"),
        ("IA / anomalies", "Complet", "Sans montants", "100 dossiers/semaine, sans montants"),
        ("Flux / SFTP / santé", "Complet", "Complet", "Interdit"),
        ("Utilisateurs / clés / exports", "Complet", "Interdit", "Interdit"),
        ("Audit des connexions", "Complet", "Technique limité", "Interdit"),
    ]
    for row in rows:
        access.append(row)
    style_sheet(access, [34, 24, 30, 38])
    security = wb.create_sheet("Sécurité pilote")
    security.append(["Couche", "Admin fondateur", "CTO", "Collaborateurs", "Statut"])
    for row in [
        ("Mot de passe", "16+ caractères", "16+ caractères", "16+ caractères", "Actif à renforcer"),
        ("TOTP", "Obligatoire", "Obligatoire", "Recommandé", "À implémenter"),
        ("Clé FIDO2/WebAuthn", "Deux clés enregistrées", "Optionnelle", "Non", "À implémenter"),
        ("SMS OTP", "Secours seulement", "Secours seulement", "Option pilote", "À implémenter"),
        ("Session", "10 min inactivité", "10 min inactivité", "10 min inactivité", "À implémenter"),
        ("Journal d’audit", "Accès complet", "Technique limité", "Non", "À implémenter"),
        ("Géoblocage", "Risque adaptatif", "Risque adaptatif", "Risque adaptatif", "À valider"),
    ]:
        security.append(row)
    style_sheet(security, [28, 32, 32, 32, 22])
    notes = wb.create_sheet("Notes gouvernance")
    notes.append(["Principe", "Exigence"])
    for row in [
        ("Minimisation", "Aucun nom client ; uniquement agrégats, HASH et zones."),
        ("Scoring", "Chaque bloc est plafonné ; score total plafonné à 100."),
        ("Décision humaine", "Le score priorise un contrôle ; il ne prouve pas une fraude et ne coupe pas automatiquement."),
        ("Article 12", "MÉROÉ fournit à la SEEG un accès lecture seule au dashboard dédié. MÉROÉ conserve l’administration exclusive de son dashboard financier et gère les niveaux d’accès."),
        ("2FA", "TOTP/FIDO2 à implémenter avant de présenter la sécurité comme opérationnelle."),
    ]:
        notes.append(row)
    style_sheet(notes, [28, 105])
    wb.save(XLSX)


def add_background(slide):
    if BACKGROUND.exists():
        slide.shapes.add_picture(str(BACKGROUND), 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor.from_string(NAVY)


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(.65), Inches(.45), Inches(12), Inches(.8))
    p = box.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor.from_string(WHITE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(.68), Inches(1.18), Inches(11.5), Inches(.45))
        p = sub.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(12); p.font.color.rgb = RGBColor.from_string(GREEN)


def add_bullets(slide, items, x=.8, y=1.8, w=11.7, h=4.8):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = shape.text_frame; tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph(); p.text = item; p.font.size = Pt(20); p.font.color.rgb = RGBColor.from_string(WHITE); p.space_after = Pt(16)


def build_presentation():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slides = [
        ("MÉROÉ — Proposition DG V5", "Pilote 3 mois · Contrôle, recouvrement et gouvernance", ["Une plateforme, deux rails : PROTEC + IA anomalies", "Une règle : la donnée utile au bon rôle, jamais plus"]),
        ("Le problème à résoudre", "Réduire les pertes et accélérer l’encaissement", ["Prévenir les coupures et améliorer la recharge", "Détecter les anomalies explicables", "Prouver le recouvrement avant facturation"]),
        ("Architecture du pilote", "Cinq flux, une traçabilité de bout en bout", ["Airtel J+1 · Guichets J+7 · SMS interne", "IA interne · Retour SEEG J+30", "Aucun nom client dans les dashboards"]),
        ("Trois niveaux d’accès", "Toi = coffre-fort · CTO = salle des machines · Équipe = salle de travail", ["Admin fondateur : finance, paramétrage, audit et exports", "CTO : technique et opérationnel, zéro argent", "Collaborateurs : lecture limitée, mission only"]),
        ("Sécurité du pilote", "Défense en profondeur, sans promesse excessive", ["TOTP obligatoire pour Admin et CTO", "Deux clés FIDO2/WebAuthn pour l’Admin", "Sessions courtes, journal d’audit et alertes", "SMS OTP seulement comme solution de secours"]),
        ("Barème unifié IA MÉROÉ — 100 points", "5 blocs plafonnés · décision explicable et auditable", ["B1 Données abonné : 20", "B2 Historique 6 mois : 30", "B3 Recharge / solde : 20", "B4 Technique compteur : 20", "B5 Géolocalisation : 10", "80–100 : contrôle urgent ; validation SEEG avant toute action"]),
        ("Tableau de décision", "Le score priorise ; il ne condamne pas", ["0–30 VERT : RAS", "31–59 ORANGE : watchlist J+7", "60–79 ROUGE : contrôle terrain prioritaire", "80–100 NOIR : contrôle urgent et décision humaine"]),
        ("Décision proposée", "Lancer le pilote 3 mois avec gouvernance validée", ["Valider les cinq flux et les responsables", "Valider la matrice d’accès et l’article 12", "Implémenter TOTP/FIDO2 avant données réelles", "Mesurer anomalies, transformation, recouvrement et valeur PROTEC"]),
    ]
    for title, subtitle, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6]); add_background(slide); add_title(slide, title, subtitle); add_bullets(slide, bullets)
    prs.save(PPTX)


if __name__ == "__main__":
    DOCS.mkdir(exist_ok=True)
    build_workbook(); build_presentation()
    print(XLSX); print(PPTX)
